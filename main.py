from fastapi import FastAPI, File, UploadFile
import fitz  # PyMuPDF for PDFs
from pptx import Presentation
import google.generativeai as genai
from collections import Counter

# Configure Gemini API
genai.configure(api_key="AIzaSyBP6SNk8wP7Gkas8oZLwefXEjlFaglQ5hc")

app = FastAPI()

# Extract text from PDFs
def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = "\n".join([page.get_text() for page in doc])
    return text

# Extract text from PPTs
def extract_text_from_ppt(file):
    prs = Presentation(file)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    return text

# Summarize text using Gemini
def summarize_text(text):
    model = genai.GenerativeModel("gemini-pro")
    response = model.generate_content(f"Summarize this: {text}")
    return response.text

# Find most repeated PYQs
def find_common_questions(text):
    questions = [line for line in text.split("\n") if "?" in line]
    most_common = Counter(questions).most_common(5)  # Top 5 repeated questions
    return [q[0] for q in most_common]

# API to upload and summarize file
@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    ext = file.filename.split(".")[-1]

    # Extract text
    if ext == "pdf":
        text = extract_text_from_pdf(file)
    elif ext == "pptx":
        text = extract_text_from_ppt(file)
    else:
        return {"error": "Only PDF and PPTX files are supported."}

    # Get summary and repeated questions
    summary = summarize_text(text)
    common_questions = find_common_questions(text)

    return {"summary": summary, "important_questions": common_questions}

# Run server
# uvicorn main:app --reload
